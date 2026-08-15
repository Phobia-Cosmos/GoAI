"""Apply the current public product language to the visual submission deck."""
from pathlib import Path
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "artifact" / "企业经营决策Agent_项目书.pptx"
TARGET = ROOT / "artifact" / "StratPilot_项目书_初赛版.pptx"


def replace_shape_text(shape, replacements):
    if not hasattr(shape, "text") or not shape.text:
        return
    text = shape.text
    updated = text
    for old, new in replacements:
        updated = updated.replace(old, new)
    if updated != text:
        shape.text = updated


def main():
    presentation = Presentation(SOURCE)
    global_replacements = [
        ("GOAI · 项目提交材料", "StratPilot · 初赛提交材料"),
        ("企业经营决策 / Agent", "StratPilot / 决策智能体"),
        ("企业经营决策 Agent", "StratPilot 决策智能体"),
        ("企业经营决策\nAgent", "StratPilot\n决策智能体"),
        ("GoAI", "StratPilot"),
        ("管理审批", "人工确认/提交"),
        ("五级来源标签", "五类来源标签"),
        ("真实企业系统接入、行业参数校准与生产验证仍需试点", "真实企业接入与行业参数校准仍需授权后试点"),
    ]
    slide_replacements = {
        3: [("应用与验证", "可点击平台与每场规则生成")],
        6: [("六步经营闭环，把数据接入、方案推演和管理审批连在一起", "六步闭环，把每场规则、方案推演、人工确认和反馈连在一起")],
        7: [("Agent 交互与解释层", "我方决策智能体交互与解释层"), ("同一经营状态内核同时服务历史复盘与未来方案推演", "模拟器、我方 Agent 和评估器通过标准接口解耦；模拟器独立负责财务真值")],
        8: [("五项工程设计，让 Agent 的建议可以验证和复核", "五项工程设计，让人机建议可以解释、验证和复盘"), ("五级来源标签", "五类来源标签")],
        9: [("典型场景：在月度经营会上比较三套可执行计划", "Demo 场景：部分可观测条件下的人机协同季度决策"), ("面向经营计划部、销售、供应链、生产与财务共同参与的滚动决策。", "只服务我方一家企业；输入本企业私有状态、公开订单和规则通知，输出带理由的候选动作与风险复盘。"), ("经营决策 Agent", "六专业 Agent 联合建议"), ("所有方案使用同一业务口径", "最终动作由人工确认，环境独立结算获单、违约和破产")],
        10: [("当前原型已覆盖经营计算闭环，企业级接入仍需试点校准", "当前 Demo 已覆盖规则生成、双阶段决策和独立结算，真实接入仍需授权试点"), ("当前定位为可运行的决策原型。", "当前定位为可运行的模拟决策原型；内部三个订单随机样本达到历史基准的相近数量级，不能表述为官方裁判精确复刻。")],
        14: [("开源建设分三阶段推进，先补治理，再稳接口，最后做社区共建", "分层开放与脱敏发布分三阶段推进，先做材料，再稳接口，最后按授权开放"), ("开源建设", "发布建设"), ("开源许可证", "发布许可证与脱敏清单"), ("社区共建", "授权后的接口与样例共建"), ("代码仓库已经公开。材料制作时仓库根目录未见 LICENSE 文件，因此许可证仍列为近期需要明确的事项。", "初赛不公开内部规则提取数据、真实企业轨迹和完整源码；另建无旧历史的脱敏发布仓库，必要时只提供编译后的可运行文件。")],
        15: [("OPEN SOURCE", "CONTROLLED RELEASE"), ("github.com/ / Phobia-Cosmos/GoAI", "github.com/Phobia-Cosmos/StratPilot"), ("企业经营决策 Agent", "StratPilot"), ("经营决策\nAgent", "StratPilot")],
    }
    for index, slide in enumerate(presentation.slides, 1):
        replacements = global_replacements + slide_replacements.get(index, [])
        for shape in slide.shapes:
            replace_shape_text(shape, replacements)
    presentation.save(TARGET)
    print(TARGET)


if __name__ == "__main__":
    main()
